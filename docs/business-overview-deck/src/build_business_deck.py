from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
DECK_DIR = ROOT / "docs" / "business-overview-deck"
OUT = DECK_DIR / "output"
PREVIEWS = DECK_DIR / "previews"
OUT.mkdir(parents=True, exist_ok=True)
PREVIEWS.mkdir(parents=True, exist_ok=True)

W, H = 1600, 900
BG = "#FFFFFF"
INK = "#17212F"
TEXT = "#2E3B4B"
MUTED = "#5D6978"
LIGHT = "#F5F7FA"
LIGHT_BLUE = "#EEF5FC"
LINE = "#CDD6E0"
BLUE = "#1D5FA7"
BLUE_DARK = "#174C86"
ORANGE = "#C57918"
GREEN = "#227C55"
RED = "#B83A3A"

FONT = "C:/Windows/Fonts/msyh.ttc"
FONT_BOLD = "C:/Windows/Fonts/msyhbd.ttc"


def font(size, bold=False):
    return ImageFont.truetype(FONT_BOLD if bold else FONT, size)


def text_lines(draw, text, f, max_width):
    lines = []
    for para in text.split("\n"):
        if para == "":
            lines.append("")
            continue
        cur = ""
        for ch in para:
            trial = cur + ch
            if draw.textlength(trial, font=f) <= max_width or not cur:
                cur = trial
            else:
                lines.append(cur)
                cur = ch
        if cur:
            lines.append(cur)
    return lines


def draw_text(draw, xy, text, size=28, fill=TEXT, bold=False, max_width=None, line_gap=8):
    x, y = xy
    f = font(size, bold)
    if max_width is None:
        draw.text((x, y), text, font=f, fill=fill)
        return draw.textbbox((x, y), text, font=f)[3]
    for line in text_lines(draw, text, f, max_width):
        draw.text((x, y), line, font=f, fill=fill)
        y += size + line_gap
    return y


def rect(draw, box, fill=BG, outline=LINE, width=2):
    draw.rectangle(box, fill=fill, outline=outline, width=width)


def line(draw, xy, fill=LINE, width=2):
    draw.line(xy, fill=fill, width=width)


def arrow(draw, x1, y1, x2, y2, color=BLUE):
    draw.line((x1, y1, x2, y2), fill=color, width=4)
    draw.polygon([(x2, y2), (x2 - 15, y2 - 8), (x2 - 15, y2 + 8)], fill=color)


def canvas():
    img = Image.new("RGB", (W, H), BG)
    d = ImageDraw.Draw(img)
    d.rectangle((0, 0, W, 10), fill=BLUE_DARK)
    return img


def header(draw, title, judgement):
    draw_text(draw, (70, 42), title, 42, INK, True, 1450)
    draw.rectangle((70, 110, 146, 116), fill=BLUE)
    draw_text(draw, (70, 132), f"核心判断：{judgement}", 24, TEXT, False, 1430, 6)
    line(draw, (70, 178, 1530, 178), LINE, 2)


def conclusion(draw, text):
    y = 812
    draw.rectangle((70, y - 8, 78, y + 52), fill=ORANGE)
    draw_text(draw, (96, y), "关键结论：", 25, ORANGE, True)
    draw_text(draw, (225, y), text, 25, TEXT, True, 1260, 6)


def footer(draw, idx):
    draw_text(draw, (70, 862), f"Bill Verification AI Pre-audit | {idx}/8", 17, "#8A96A3")


def save_slide(img, idx):
    path = PREVIEWS / f"slide_{idx:02d}.png"
    img.save(path)
    return path


def table(draw, x, y, col_widths, row_heights, rows, header_fill=LIGHT_BLUE):
    cur_y = y
    for r, row_h in enumerate(row_heights):
        cur_x = x
        for c, col_w in enumerate(col_widths):
            fill = header_fill if r == 0 else (LIGHT if r % 2 == 0 else BG)
            rect(draw, (cur_x, cur_y, cur_x + col_w, cur_y + row_h), fill=fill, outline=LINE, width=2)
            item = rows[r][c]
            color = BLUE if r == 0 else TEXT
            bold = r == 0 or (c == 0 and r > 0)
            draw_text(draw, (cur_x + 18, cur_y + 16), item, 22 if r == 0 else 21, color, bold, col_w - 36, 6)
            cur_x += col_w
        cur_y += row_h


def s1():
    img = canvas()
    d = ImageDraw.Draw(img)
    header(d, "构建模板字段驱动的权签票据 AI 预审能力", "业务先定义当前模板要核验哪些系统字段，AI 按字段清单定向提取，系统再用别名规则补映射并提示风险。")

    y = 245
    labels = [
        ("选择核验模板", "带出本次待检查系统字段"),
        ("AI 定向提取", "输出 document_items / extracted_fields"),
        ("别名补映射", "模板别名规则兜底生成字段映射"),
        ("函数包核验", "系统指令 vs 票面字段，输出风险等级"),
        ("权签人确认", "聚焦风险并沉淀优化反馈"),
    ]
    x = 75
    w = 250
    for i, (title, body) in enumerate(labels):
        rect(d, (x, y, x + w, y + 165), fill=LIGHT_BLUE if i in (1, 2) else BG)
        draw_text(d, (x + 18, y + 25), f"{i + 1}. {title}", 25, BLUE, True, w - 36)
        draw_text(d, (x + 18, y + 78), body, 20, TEXT, False, w - 36, 5)
        if i < len(labels) - 1:
            arrow(d, x + w + 8, y + 83, x + w + 55, y + 83, BLUE)
        x += w + 62

    rows = [
        ["建设对象", "核心能力", "管理价值"],
        ["权签预审场景", "模板字段清单、票面原文、系统映射、风险结果", "让权签人先看风险项，而不是逐字肉眼遍历票面"],
        ["模板运营机制", "模板新增/复制、字段增删改、别名、位置提示、反馈样本", "让字段叫法和模板差异通过配置持续迭代"],
    ]
    table(d, 130, 500, [270, 520, 520], [54, 92, 92], rows)
    conclusion(d, "最新方案把“该提什么”前移到模板配置，把“怎么映射”落实为 Prompt 辅助 + 后端规则兜底。")
    footer(d, 1)
    return img


def s2():
    img = canvas()
    d = ImageDraw.Draw(img)
    header(d, "从人工逐项核对升级为 AI 辅助风险初筛", "现有权签核对面对模板多、字段多、语言多和人员分散的复合压力，需要用智能化能力先完成一轮风险聚焦。")

    rows = [
        ["压力来源", "现状表现", "对权签人的影响", "AI 预审切入点"],
        ["票据模板多", "全球 200+ 银行模板，版式、语言和字段叫法差异明显", "难以依赖统一肉眼经验快速判断", "按模板沉淀别名、位置和抽取提示"],
        ["系统字段多", "MAP/AP 指令有几十个结构化字段，票面只出现子集", "需要判断哪些字段应参与本次核验", "先抽取票面原文，再映射到系统字段"],
        ["兼职权签多", "600+ 权签人中大量兼职权签，职位层级较高", "使用体验对准确性和解释性要求高", "突出风险项，通过项收起，保留证据链"],
        ["迭代周期长", "办公系统版本周期约两个月，频繁发版成本高", "模板问题难以及时修正", "通过配置和反馈样本形成运营闭环"],
    ]
    table(d, 70, 220, [210, 430, 380, 390], [58, 95, 95, 95, 95], rows)

    d.rectangle((105, 716, 1495, 760), fill=LIGHT_BLUE, outline=LINE)
    draw_text(d, (130, 726), "评估口径建议：准确率、召回率、风险采纳率、反馈闭环率、权签人使用满意度。", 24, BLUE_DARK, True, 1320)
    conclusion(d, "第一阶段应以“高风险不漏报、误报可解释”为目标，避免把 Demo 做成只看模型炫技的识别工具。")
    footer(d, 2)
    return img


def s3():
    img = canvas()
    d = ImageDraw.Draw(img)
    header(d, "将 AI 预审拆成字段配置、模型提取和规则补映射", "模型负责看图和提取候选，系统负责限定字段范围、补足别名映射和保留可解释证据链。")

    rows = [
        ["层级", "输入", "输出", "控制方式", "关键风险"],
        ["1. 模板字段配置", "银行/国家/单据模板 + MAP 字段", "待检查字段、含义、别名、位置提示", "只有发布模板才启动 AI", "未配置模板无法进入预审"],
        ["2. 模型定向提取", "票面图片 + 当前模板字段清单", "双层 JSON 结果", "Prompt 要求保留原文并限制字段集合", "模型只返回原文但漏掉字段映射"],
        ["3. 后端补映射", "document_items + 模板别名", "mapped_field、系统来源字段、mapping_source", "别名规则按模板范围确定性生效", "别名过宽导致误映射"],
        ["4. 规则核验", "系统支付指令 + 标准字段", "风险项、通过项、证据和置信度", "MAP/交易系统函数包负责最终比对", "比对规则归属不清"],
    ]
    table(d, 60, 220, [190, 310, 350, 350, 250], [58, 105, 105, 105, 105], rows)
    conclusion(d, "当前设计要求 JSON 同时回答三件事：票面写了什么、对应哪个系统字段、这个映射来自模型还是模板别名规则。")
    footer(d, 3)
    return img


def s4():
    img = canvas()
    d = ImageDraw.Draw(img)
    header(d, "核验页从“选择模板”开始逐字段输出预审结果", "权签人进入页面先看到模板字段清单，AI 预审后每个待检查系统字段都有提取结果或风险提示。")

    y = 260
    x = 90
    steps = [
        ("选择模板", "展示待检查系统字段\n样例任务可自动带出模板"),
        ("输入任务", "系统支付指令 + 票面文件\n并列预览"),
        ("AI 逐字段提取", "字段有值则展示证据\n未提取则标记风险"),
        ("提交优化反馈", "例如新增别名\n入账行 -> beneficiary_bank"),
        ("重新预审验证", "重新真实提取与核验\n确认配置生效"),
    ]
    for i, (title, body) in enumerate(steps):
        rect(d, (x, y, x + 245, y + 170), fill=LIGHT if i % 2 == 0 else BG)
        draw_text(d, (x + 20, y + 24), f"{i + 1}", 31, ORANGE if i == 3 else BLUE, True)
        draw_text(d, (x + 64, y + 30), title, 24, INK, True, 160)
        draw_text(d, (x + 20, y + 82), body, 19, TEXT, False, 205, 4)
        if i < len(steps) - 1:
            arrow(d, x + 255, y + 86, x + 305, y + 86, BLUE)
        x += 292

    rows = [
        ["演示用例", "初始问题", "人工动作", "验证结果"],
        ["入账行别名漏识别", "原文有“入账行”，但结构化结果没有收款方银行字段", "在当前模板下新增别名映射并保存", "后端补出 mapped_field，再次预审后风险消失或降级"],
    ]
    table(d, 140, 540, [270, 520, 340, 300], [54, 96], rows)
    conclusion(d, "业务用户看到的是逐字段风险闭环，产品侧沉淀的是模板字段、别名规则和可复用反馈样本。")
    footer(d, 4)
    return img


def s5():
    img = canvas()
    d = ImageDraw.Draw(img)
    header(d, "别名配置采用“Prompt 辅助 + 后端补映射”双通道生效", "业务改一个字段叫法后，既能帮助模型理解票面，也能让后端在模型漏映射时确定性补齐字段。")

    rows = [
        ["配置项", "进入 Prompt 的作用", "进入后端规则的作用", "边界控制"],
        ["字段别名", "告诉模型票面可能叫法，例如入账行", "raw_key 命中别名时补系统字段映射", "限定在当前模板字段内"],
        ["业务含义", "帮助模型区分相似字段，例如收款银行 vs 中间行", "不直接做规则匹配，只作为解释和审阅信息", "避免自然语言描述变成泛化规则"],
        ["位置提示", "帮助模型在版面上找字段区域", "不参与字段映射，只辅助证据判断", "避免位置配置污染其他模板"],
        ["反馈优化", "下次调用模型时带入新叫法", "保存为模板别名或字段提示，支持回滚", "按模板版本发布"],
    ]
    table(d, 60, 220, [190, 420, 430, 410], [58, 95, 95, 95, 95], rows)
    conclusion(d, "配置体系不再只是“把词塞进 Prompt”，而是让业务配置同时约束模型输入和系统后处理。")
    footer(d, 5)
    return img


def s6():
    img = canvas()
    d = ImageDraw.Draw(img)
    header(d, "将高频动作拆分为四类页面，降低业务理解成本", "模型配置、付款核验、模板调优和反馈样本应分工明确，避免把一次性设置和日常核验混在同一工作台。")

    rows = [
        ["页面", "主要用户", "核心动作", "展示重点", "当前 Demo 状态"],
        ["付款核验", "权签人", "选择模板、查看字段清单、上传票面、启动 AI 预审、提交优化", "每个模板字段都有提取结果或风险提示", "主流程已打通"],
        ["模板调优", "业务配置人员", "新增/复制模板，增删改字段，维护别名、位置提示、提取要求", "少技术术语，按模板和字段组织", "已支持基础管理"],
        ["反馈样本", "运营 / AI 产品", "查看纠错、忽略、确认不一致、提交优化记录", "高频问题聚合，转规则或转样本", "已展示记录"],
        ["系统设置", "管理员", "配置模型 URL、模型名、API Key、文本/图片诊断", "一次性配置，不干扰日常核验", "已可配置和测试"],
    ]
    table(d, 55, 218, [190, 180, 430, 390, 260], [58, 100, 100, 100, 100], rows)
    conclusion(d, "面向演示时应突出付款核验主流程，模板调优和系统设置作为支撑能力放在后面说明。")
    footer(d, 6)
    return img


def s7():
    img = canvas()
    d = ImageDraw.Draw(img)
    header(d, "沉淀可修改测试资产，支撑业务持续验证配置能力", "真实票据敏感不可外传，因此需要用可编辑 Word 和合成截图构建可复用、可扩展的测试样本体系。")

    rows = [
        ["资产类型", "文件 / 样例", "覆盖重点", "使用方式"],
        ["中文支票 Word", "cn_check_alias_feedback_case", "付款人、收款人、入账行、金额、币种、不可转让", "业务修改字段名和值后截图上传"],
        ["英文电汇 Word", "en_tt_payment_application_case", "Ordering Customer、Bank、SWIFT", "验证英文和跨境字段映射"],
        ["预置样例 JSON", "中文通过、金额差异、账号差异、别名漏识别", "覆盖风险核验和反馈闭环", "用于 Demo 快速切换与调试"],
        ["模型诊断脚本", "test_model_api.py / image variants", "文本链路、图片链路、base64 image_url 兼容性", "定位公司网关和模型调用差异"],
    ]
    table(d, 70, 220, [230, 390, 460, 350], [58, 100, 100, 100, 100], rows)
    d.rectangle((120, 700, 1480, 748), fill=LIGHT_BLUE, outline=LINE)
    draw_text(d, (145, 711), "建议指标位：样例覆盖数、字段覆盖率、模板覆盖率、反馈闭环率、模型结构化成功率。", 24, BLUE_DARK, True, 1300)
    conclusion(d, "测试资产需要像规则一样持续沉淀，后续才能支撑不同银行模板和不同语言场景的能力评估。")
    footer(d, 7)
    return img


def s8():
    img = canvas()
    d = ImageDraw.Draw(img)
    header(d, "以 MVP 验证带动模板运营和规模化推广", "当前 Demo 已验证模板选择、真实模型提取、别名补映射、风险核验和反馈闭环，下一步应接入 MAP 配置与函数包。")

    rows = [
        ["阶段", "建设目标", "关键动作", "衡量指标"],
        ["阶段一：Demo 验证", "打通模板字段、模型提取、补映射、风险核验和反馈闭环", "完善样例、刷新页面旅程、沉淀文档和汇报材料", "结构化成功率、演示通过率"],
        ["阶段二：试点模板", "选取高频国家/银行模板进行小范围试点", "建立模板选择、别名范围、特殊条款配置和反馈运营机制", "风险召回率、误报率、反馈闭环率"],
        ["阶段三：MAP 接入", "由 MAP 管理系统字段、模板发布和比对函数包", "打通任务输入、票面上传、结果回写和权限边界", "流程耗时、采纳率、问题定位效率"],
        ["阶段四：规模运营", "形成全球模板资产与规则运营体系", "建立样本池、规则版本、回滚机制和效果报表", "模板覆盖率、规则沉淀量、用户满意度"],
    ]
    table(d, 65, 220, [230, 390, 540, 290], [58, 102, 102, 102, 102], rows)
    conclusion(d, "后续重点不是单次模型效果展示，而是形成“模板资产 + 规则资产 + 反馈样本 + 效果指标”的持续运营机制。")
    footer(d, 8)
    return img


slides = [s1(), s2(), s3(), s4(), s5(), s6(), s7(), s8()]
preview_paths = [save_slide(img, i + 1) for i, img in enumerate(slides)]


def add_textbox(slide, x, y, w, h, text_value, size=8, color="#7A8795"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text_value
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor.from_string(color.strip("#"))
    return box


prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
blank = prs.slide_layouts[6]
for idx, png in enumerate(preview_paths, start=1):
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(png), 0, 0, width=Inches(16), height=Inches(9))
    add_textbox(slide, 0.15, 8.55, 6, 0.3, f"权签票据一致性 AI 预审 - 第 {idx} 页")

pptx_path = OUT / "bill_verification_business_overview.pptx"
prs.save(pptx_path)

montage = Image.new("RGB", (W * 2, H * 4), "#FFFFFF")
for i, img in enumerate(slides):
    montage.paste(img.resize((W, H)), ((i % 2) * W, (i // 2) * H))
montage_path = PREVIEWS / "montage.png"
montage.save(montage_path)

print(pptx_path)
for p in preview_paths:
    print(p)
print(montage_path)
