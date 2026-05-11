from __future__ import annotations

import json
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
DECK_DIR = ROOT / "docs" / "business-overview-deck"
OUT = DECK_DIR / "output"
OUT.mkdir(parents=True, exist_ok=True)

PPTX_PATH = OUT / "bill_verification_business_overview_editable.pptx"
SUMMARY_PATH = OUT / "bill_verification_business_overview_editable_summary.json"

W, H = 16, 9
BG = "FFFFFF"
INK = "17212F"
TEXT = "2E3B4B"
MUTED = "5D6978"
LIGHT = "F5F7FA"
LIGHT_BLUE = "EEF5FC"
LINE = "CDD6E0"
BLUE = "1D5FA7"
BLUE_DARK = "174C86"
ORANGE = "C57918"


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_typeface(run, name: str = "Microsoft YaHei") -> None:
    run.font.name = name
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = r_pr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            r_pr.append(node)
        node.set("typeface", name)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 14,
    color: str = TEXT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(0.03)
    box.text_frame.margin_right = Inches(0.03)
    box.text_frame.margin_top = Inches(0.02)
    box.text_frame.margin_bottom = Inches(0.02)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    set_typeface(run)
    return box


def add_rect(slide, x: float, y: float, w: float, h: float, fill: str = "FFFFFF", line: str = LINE, width: float = 1):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(width)
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color: str = LINE, width: float = 1.2):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    return shape


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color: str = BLUE):
    shape = add_line(slide, x1, y1, x2, y2, color, 2)
    shape.line.end_arrowhead = True
    return shape


def add_header(slide, title: str, judgement: str):
    add_rect(slide, 0, 0, W, 0.1, BLUE_DARK, BLUE_DARK, 0)
    add_text(slide, title, 0.55, 0.42, 14.5, 0.48, 22, INK, True)
    add_rect(slide, 0.55, 1.07, 0.58, 0.05, BLUE, BLUE, 0)
    add_text(slide, f"核心判断：{judgement}", 0.55, 1.34, 14.7, 0.34, 12.2, TEXT)
    add_line(slide, 0.55, 1.76, 15.45, 1.76, LINE, 1)


def add_conclusion(slide, text: str):
    add_rect(slide, 0.55, 8.15, 0.07, 0.55, ORANGE, ORANGE, 0)
    add_text(slide, "关键结论：", 0.72, 8.2, 1.15, 0.35, 12.5, ORANGE, True)
    add_text(slide, text, 1.75, 8.2, 13.4, 0.42, 12.5, INK, True)


def add_footer(slide, index: int):
    add_text(slide, f"Bill Verification AI Pre-audit | {index}/9", 0.55, 8.72, 4.3, 0.22, 8.5, "8A96A3")


def add_table(
    slide,
    rows: list[list[str]],
    x: float,
    y: float,
    widths: list[float],
    row_h: float,
    font_size: float = 10.5,
    header_size: float = 11.3,
):
    for r, row in enumerate(rows):
        cursor = x
        fill = LIGHT_BLUE if r == 0 else (LIGHT if r % 2 == 0 else BG)
        for c, text in enumerate(row):
            add_rect(slide, cursor, y + r * row_h, widths[c], row_h, fill, LINE, 1)
            add_text(
                slide,
                text,
                cursor + 0.1,
                y + r * row_h + 0.1,
                widths[c] - 0.18,
                row_h - 0.12,
                header_size if r == 0 else font_size,
                BLUE if r == 0 else INK,
                r == 0 or c == 0,
            )
            cursor += widths[c]


def add_flow(slide, labels: list[tuple[str, str]], x: float, y: float, w: float, h: float, gap: float):
    for i, (title, body) in enumerate(labels):
        add_rect(slide, x, y, w, h, LIGHT_BLUE if i in (1, 2) else BG, LINE, 1)
        add_text(slide, f"{i + 1}. {title}", x + 0.12, y + 0.25, w - 0.24, 0.28, 12.3, BLUE, True)
        add_text(slide, body, x + 0.12, y + 0.75, w - 0.24, h - 0.82, 9.5, TEXT)
        if i < len(labels) - 1:
            add_arrow(slide, x + w + 0.08, y + h / 2, x + w + gap - 0.08, y + h / 2)
        x += w + gap


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def s1(prs):
    slide = blank_slide(prs)
    add_header(slide, "构建模板字段驱动的权签票据 AI 预审能力", "业务先定义当前模板要核验哪些系统字段，AI 按字段清单定向提取，系统再用别名规则补映射并提示风险。")
    add_flow(
        slide,
        [
            ("选择核验模板", "带出本次待检查系统字段"),
            ("AI 定向提取", "输出 document_items / extracted_fields"),
            ("别名补映射", "模板别名规则兜底生成字段映射"),
            ("函数包核验", "系统指令 vs 票面字段，输出风险等级"),
            ("权签人确认", "聚焦风险并沉淀优化反馈"),
        ],
        0.6,
        2.45,
        1.55,
        1.55,
        0.4,
    )
    add_table(
        slide,
        [
            ["建设对象", "核心能力", "管理价值"],
            ["权签预审场景", "模板字段清单、票面原文、系统映射、风险结果", "让权签人先看风险项，而不是逐字肉眼遍历票面"],
            ["模板运营机制", "模板新增/复制、字段增删改、票面别名、给 AI 的输入、反馈样本", "让字段叫法和模板差异通过配置持续迭代"],
        ],
        1.1,
        5.0,
        [2.6, 5.0, 5.0],
        0.62,
        10.4,
    )
    add_conclusion(slide, "最新方案把“该提什么”前移到模板配置，把“怎么映射”落实为 Prompt 辅助 + 后端规则兜底。")
    add_footer(slide, 1)


def s2(prs):
    slide = blank_slide(prs)
    add_header(slide, "从人工逐项核对升级为 AI 辅助风险初筛", "现有权签核对面对模板多、字段多、语言多和人员分散的复合压力，需要用智能化能力先完成一轮风险聚焦。")
    add_table(
        slide,
        [
            ["压力来源", "现状表现", "对权签人的影响", "AI 预审切入点"],
            ["票据模板多", "全球 200+ 银行模板，版式、语言和字段叫法差异明显", "难以依赖统一肉眼经验快速判断", "按模板沉淀别名、位置和抽取提示"],
            ["系统字段多", "MAP/AP 指令有几十个结构化字段，票面只出现子集", "需要判断哪些字段应参与本次核验", "先抽取票面原文，再映射到系统字段"],
            ["兼职权签多", "600+ 权签人中大量兼职权签，职位层级较高", "使用体验对准确性和解释性要求高", "突出风险项，通过项收起，保留证据链"],
            ["迭代周期长", "办公系统版本周期约两个月，频繁发版成本高", "模板问题难以及时修正", "通过配置和反馈样本形成运营闭环"],
        ],
        0.55,
        2.15,
        [1.7, 4.1, 3.6, 3.6],
        0.88,
        9.5,
    )
    add_rect(slide, 0.85, 7.18, 14.1, 0.36, LIGHT_BLUE, LINE, 1)
    add_text(slide, "评估口径建议：准确率、召回率、风险采纳率、反馈闭环率、权签人使用满意度。", 1.0, 7.25, 13.5, 0.24, 11.8, BLUE_DARK, True)
    add_conclusion(slide, "第一阶段应以“高风险不漏报、误报可解释”为目标，避免把 Demo 做成只看模型炫技的识别工具。")
    add_footer(slide, 2)


def s3(prs):
    slide = blank_slide(prs)
    add_header(slide, "将 AI 预审拆成字段配置、模型提取和规则补映射", "模型负责看图和提取候选，系统负责限定字段范围、补足别名映射和保留可解释证据链。")
    add_table(
        slide,
        [
            ["层级", "输入", "输出", "控制方式", "关键风险"],
            ["1. 模板字段配置", "银行/国家/单据模板 + MAP 字段", "待检查字段、系统来源、票面别名、给 AI 的输入", "只有发布模板才启动 AI", "未配置模板无法进入预审"],
            ["2. 模型定向提取", "票面图片 + 当前模板字段清单", "双层 JSON 结果", "Prompt 要求保留原文并限制字段集合", "模型只返回原文但漏掉字段映射"],
            ["3. 后端补映射", "document_items + 模板别名", "mapped_field、系统来源字段、mapping_source", "别名规则按模板范围确定性生效", "别名过宽导致误映射"],
            ["4. 规则核验", "系统支付指令 + 标准字段", "风险项、通过项、证据和置信度", "MAP/交易系统函数包负责最终比对", "比对规则归属不清"],
        ],
        0.45,
        2.15,
        [1.55, 3.0, 3.35, 3.35, 2.35],
        0.9,
        9.2,
    )
    add_conclusion(slide, "当前设计要求 JSON 同时回答三件事：票面写了什么、对应哪个系统字段、这个映射来自模型还是模板别名规则。")
    add_footer(slide, 3)


def s4(prs):
    slide = blank_slide(prs)
    add_header(slide, "完整作业流按输入输出拆清楚", "每个环节都要明确输入、处理、输出和责任系统，避免把模型提取、字段映射和最终核验混成一个黑盒。")
    add_table(
        slide,
        [
            ["环节", "输入", "处理动作", "输出", "责任边界"],
            ["1. 模板选择", "付款任务、银行、国家、单据类型", "系统带出或人工选择模板", "template_id、待检查字段清单", "MAP/AP 或交易系统"],
            ["2. 字段配置", "模板字段、系统来源、票面别名、给 AI 的输入", "发布模板字段配置", "可调用 AI 的字段清单", "业务配置 + MAP 管控"],
            ["3. AI 预审", "票面图片、template_id、字段清单", "模型定向提取票面信息", "document_items / extracted_fields", "AI 产品"],
            ["4. 补映射", "document_items、模板别名", "raw_key 命中别名时补字段映射", "mapped_field、mapping_source", "AI 后端规则"],
            ["5. 函数包核验", "系统支付指令、标准字段、票面值", "金额、日期、账号、名称等规则比对", "风险项、通过项、证据", "MAP/交易系统"],
            ["6. 人工反馈", "风险项、证据、字段配置", "确认、忽略、纠错、提交优化", "反馈样本、别名或提示建议", "权签人 + 运营"],
        ],
        0.45,
        2.15,
        [1.65, 2.9, 3.15, 3.0, 2.6],
        0.66,
        8.5,
        9.5,
    )
    add_conclusion(slide, "产品上要把 AI 输出看成中间证据，不是最终判断；最终权签核验结果应由系统规则和人工确认共同闭环。")
    add_footer(slide, 4)


def s5(prs):
    slide = blank_slide(prs)
    add_header(slide, "核验页从“选择模板”开始逐字段输出预审结果", "权签人进入页面先看到模板字段清单，AI 预审后每个待检查系统字段都有提取结果或风险提示。")
    add_flow(
        slide,
        [
            ("选择模板", "展示待检查系统字段\n样例任务可自动带出模板"),
            ("输入任务", "系统支付指令 + 票面文件\n并列预览"),
            ("AI 逐字段提取", "字段有值则展示证据\n未提取则标记风险"),
            ("提交优化反馈", "例如新增别名\n入账行 -> beneficiary_bank"),
            ("重新预审验证", "重新真实提取与核验\n确认配置生效"),
        ],
        0.65,
        2.65,
        1.45,
        1.55,
        0.42,
    )
    add_table(
        slide,
        [
            ["演示用例", "初始问题", "人工动作", "验证结果"],
            ["入账行别名漏识别", "原文有“入账行”，但结构化结果没有收款方银行字段", "在当前模板下新增别名映射并保存", "后端补出 mapped_field，再次预审后风险消失或降级"],
        ],
        1.0,
        5.75,
        [2.3, 5.0, 3.2, 3.1],
        0.62,
        10.0,
    )
    add_conclusion(slide, "业务用户看到的是逐字段风险闭环，产品侧沉淀的是模板字段、别名规则和可复用反馈样本。")
    add_footer(slide, 5)


def s6(prs):
    slide = blank_slide(prs)
    add_header(slide, "别名配置采用“Prompt 辅助 + 后端补映射”双通道生效", "业务改一个字段叫法后，既能帮助模型理解票面，也能让后端在模型漏映射时确定性补齐字段。")
    add_table(
        slide,
        [
            ["配置项", "进入 Prompt 的作用", "进入后端规则的作用", "边界控制"],
            ["字段别名", "告诉模型票面可能叫法，例如入账行", "raw_key 命中别名时补系统字段映射", "限定在当前模板字段内"],
            ["给 AI 的输入", "合并说明字段含义、常见位置和提取注意事项", "不直接做规则匹配，只作为模型理解提示", "避免自然语言描述变成泛化规则"],
            ["反馈优化", "下次调用模型时带入新叫法或更清晰的提取说明", "保存为模板别名或 AI 输入建议，支持回滚", "按模板版本发布"],
        ],
        0.55,
        2.25,
        [1.75, 4.1, 4.2, 3.55],
        0.92,
        9.4,
    )
    add_conclusion(slide, "配置体系不再只是“把词塞进 Prompt”，而是让业务配置同时约束模型输入和系统后处理。")
    add_footer(slide, 6)


def s7(prs):
    slide = blank_slide(prs)
    add_header(slide, "将高频动作拆分为四类页面，降低业务理解成本", "模型配置、付款核验、模板调优和反馈样本应分工明确，避免把一次性设置和日常核验混在同一工作台。")
    add_table(
        slide,
        [
            ["页面", "主要用户", "核心动作", "展示重点", "当前 Demo 状态"],
            ["付款核验", "权签人", "选择模板、查看字段清单、上传票面、启动 AI 预审、提交优化", "每个模板字段都有提取结果或风险提示", "主流程已打通"],
            ["模板调优", "业务配置人员", "新增/复制模板，增删改字段，维护票面别名和给 AI 的输入", "少技术术语，按模板和字段组织", "已支持基础管理"],
            ["反馈样本", "运营 / AI 产品", "查看纠错、忽略、确认不一致、提交优化记录", "高频问题聚合，转规则或转样本", "已展示记录"],
            ["系统设置", "管理员", "配置模型 URL、模型名、API Key、文本/图片诊断", "一次性配置，不干扰日常核验", "已可配置和测试"],
        ],
        0.45,
        2.25,
        [1.65, 1.65, 4.15, 3.75, 2.15],
        0.9,
        9.0,
    )
    add_conclusion(slide, "面向演示时应突出付款核验主流程，模板调优和系统设置作为支撑能力放在后面说明。")
    add_footer(slide, 7)


def s8(prs):
    slide = blank_slide(prs)
    add_header(slide, "沉淀可修改测试资产，支撑业务持续验证配置能力", "真实票据敏感不可外传，因此需要用可编辑 Word 和合成截图构建可复用、可扩展的测试样本体系。")
    add_table(
        slide,
        [
            ["资产类型", "文件 / 样例", "覆盖重点", "使用方式"],
            ["中文支票 Word", "cn_check_alias_feedback_case", "付款人、收款人、入账行、金额、币种、不可转让", "业务修改字段名和值后截图上传"],
            ["英文电汇 Word", "en_tt_payment_application_case", "Ordering Customer、Bank、SWIFT", "验证英文和跨境字段映射"],
            ["预置样例 JSON", "中文通过、金额差异、账号差异、别名漏识别", "覆盖风险核验和反馈闭环", "用于 Demo 快速切换与调试"],
            ["模型诊断脚本", "test_model_api.py / image variants", "文本链路、图片链路、base64 image_url 兼容性", "定位公司网关和模型调用差异"],
        ],
        0.55,
        2.25,
        [2.1, 3.7, 4.3, 3.4],
        0.9,
        9.4,
    )
    add_rect(slide, 0.95, 7.18, 13.8, 0.36, LIGHT_BLUE, LINE, 1)
    add_text(slide, "建议指标位：样例覆盖数、字段覆盖率、模板覆盖率、反馈闭环率、模型结构化成功率。", 1.1, 7.25, 13.1, 0.24, 11.5, BLUE_DARK, True)
    add_conclusion(slide, "测试资产需要像规则一样持续沉淀，后续才能支撑不同银行模板和不同语言场景的能力评估。")
    add_footer(slide, 8)


def s9(prs):
    slide = blank_slide(prs)
    add_header(slide, "以 MVP 验证带动模板运营和规模化推广", "当前 Demo 已验证模板选择、真实模型提取、别名补映射、风险核验和反馈闭环，下一步应接入 MAP 配置与函数包。")
    add_table(
        slide,
        [
            ["阶段", "建设目标", "关键动作", "衡量指标"],
            ["阶段一：Demo 验证", "打通模板字段、模型提取、补映射、风险核验和反馈闭环", "完善样例、刷新页面旅程、沉淀文档和汇报材料", "结构化成功率、演示通过率"],
            ["阶段二：试点模板", "选取高频国家/银行模板进行小范围试点", "建立模板选择、别名范围、特殊条款配置和反馈运营机制", "风险召回率、误报率、反馈闭环率"],
            ["阶段三：MAP 接入", "由 MAP 管理系统字段、模板发布和比对函数包", "打通任务输入、票面上传、结果回写和权限边界", "流程耗时、采纳率、问题定位效率"],
            ["阶段四：规模运营", "形成全球模板资产与规则运营体系", "建立样本池、规则版本、回滚机制和效果报表", "模板覆盖率、规则沉淀量、用户满意度"],
        ],
        0.55,
        2.25,
        [2.05, 3.6, 5.0, 2.75],
        0.9,
        9.4,
    )
    add_conclusion(slide, "后续重点不是单次模型效果展示，而是形成“模板资产 + 规则资产 + 反馈样本 + 效果指标”的持续运营机制。")
    add_footer(slide, 9)


def count_shapes(prs: Presentation) -> dict[str, int]:
    counts = {"slides": len(prs.slides), "text_boxes": 0, "shapes": 0, "tables": 0, "pictures": 0}
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                counts["tables"] += 1
            elif getattr(shape, "shape_type", None) == 13:
                counts["pictures"] += 1
            elif shape.has_text_frame:
                counts["text_boxes"] += 1
            else:
                counts["shapes"] += 1
    return counts


def main():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    for builder in (s1, s2, s3, s4, s5, s6, s7, s8, s9):
        builder(prs)

    prs.save(PPTX_PATH)

    reopened = Presentation(PPTX_PATH)
    summary = {
        "pptx": str(PPTX_PATH),
        "editable": True,
        "fallback_images": 0,
        "note": "All visible slide content is built from native PowerPoint text boxes, shapes, lines, and editable table-like cell groups. No full-slide image fallback is used.",
        "counts": count_shapes(reopened),
    }
    SUMMARY_PATH.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    print(PPTX_PATH)
    print(SUMMARY_PATH)
    print(json.dumps(summary["counts"], ensure_ascii=False))


if __name__ == "__main__":
    main()
